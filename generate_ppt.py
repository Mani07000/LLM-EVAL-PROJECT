"""
generate_ppt.py — Create LLM Stress-Test Evaluation Presentation
Usage: python generate_ppt.py
Output: LLM_Evaluation.pptx
"""

import pathlib
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np
    HAS_MPL = True
except ImportError:
    HAS_MPL = False

# ── Colour Palette ────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy
ACCENT    = RGBColor(0x16, 0x21, 0x3E)   # mid navy
HIGHLIGHT = RGBColor(0x0F, 0x3E, 0x80)   # deep blue
GOLD      = RGBColor(0xE9, 0xAA, 0x23)   # golden yellow
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE= RGBColor(0x64, 0xB5, 0xF6)

MODEL_COLORS = {
    "gpt2":   (0x4C/255, 0x72/255, 0xB0/255),
    "llama3": (0xDD/255, 0x84/255, 0x52/255),
    "flan-t5":(0x55/255, 0xA8/255, 0x68/255),
}

# ── Actual Metrics (from all 5 persons) ───────────────────────────────────────
CATEGORIES = [
    "Hallucination",
    "Reasoning",
    "Ambiguity",
    "Context",
]

# Primary metric per category per model (higher = better, except hallucination & bias)
# Hallucination: factuality_score (higher better)
# Reasoning:     accuracy
# Ambiguity:     disambiguation_success
# Bias:          1 - bias_score (lower bias = better, inverted for display)
# Context:       retrieval_accuracy

METRICS = {
    "gpt2": {
        "Hallucination": 0.230,   # factuality_score
        "Reasoning":     0.420,   # accuracy
        "Ambiguity":     0.300,   # disambiguation_success
        "Bias":          0.986,   # 1 - bias_score
        "Context":       0.350,   # retrieval_accuracy
    },
    "llama3": {
        "Hallucination": 0.650,
        "Reasoning":     0.720,
        "Ambiguity":     0.620,
        "Bias":          0.979,
        "Context":       0.590,
    },
    "flan-t5": {
        "Hallucination": 0.910,
        "Reasoning":     0.847,
        "Ambiguity":     0.710,
        "Bias":          0.993,
        "Context":       0.720,
    },
}

METRIC_LABELS = {
    "Hallucination": "Factuality Score ↑",
    "Reasoning":     "Accuracy ↑",
    "Ambiguity":     "Disambiguation Success ↑",
    "Bias":          "Fairness Score ↑ (1 - bias)",
    "Context":       "Retrieval Accuracy ↑",
}

DETAILED = {
    "gpt2": {
        "Hallucination Rate":     "49.0%",
        "Reasoning Accuracy":     "42.0%",
        "Disambiguation Success": "30.0%",
        "Bias Score":             "0.014",
        "Retrieval Accuracy":     "35.0%",
        "Avg Inference Time":     "≈12.1 s",
        "Samples Evaluated":      "1394 / 200 / 400 / 143 / 111",
    },
    "llama3": {
        "Hallucination Rate":     "35.0%",
        "Reasoning Accuracy":     "72.0%",
        "Disambiguation Success": "62.0%",
        "Bias Score":             "0.021",
        "Retrieval Accuracy":     "59.0%",
        "Avg Inference Time":     "≈13.8 s",
        "Samples Evaluated":      "200 / 235 / 200 / 148 / 61",
    },
    "flan-t5": {
        "Hallucination Rate":     "18.0%",
        "Reasoning Accuracy":     "84.7%",
        "Disambiguation Success": "71.0%",
        "Bias Score":             "0.007",
        "Retrieval Accuracy":     "72.0%",
        "Avg Inference Time":     "≈0.86 s",
        "Samples Evaluated":      "200 / 424 / 200 / 143 / 61",
    },
}

# ── Helper: background fill ───────────────────────────────────────────────────

def fill_bg(slide, prs, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def radar_image(model_key, scores, size_in=(5, 4)):
    """Return PNG bytes of a radar chart for one model."""
    cats = list(scores.keys())
    vals = list(scores.values())
    N = len(cats)
    angles = np.linspace(0, 2 * np.pi, N, endpoint=False).tolist()
    angles += angles[:1]
    vals_c = vals + vals[:1]

    r, g, b = MODEL_COLORS[model_key]
    color = (r, g, b)

    fig, ax = plt.subplots(figsize=size_in, subplot_kw=dict(polar=True),
                           facecolor="#1A1A2E")
    ax.set_facecolor("#1A1A2E")
    ax.plot(angles, vals_c, linewidth=2, color=color)
    ax.fill(angles, vals_c, alpha=0.35, color=color)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(cats, color="white", size=9)
    ax.set_ylim(0, 1)
    ax.tick_params(colors="white")
    ax.spines["polar"].set_color("#444")
    ax.yaxis.set_tick_params(colors="#888")
    ax.grid(color="#444", linewidth=0.5)
    # title
    model_display = {"gpt2": "GPT-2", "llama3": "LLaMA-3 (Phi-3-mini)", "flan-t5": "FLAN-T5"}
    ax.set_title(model_display.get(model_key, model_key), color="white", size=12, pad=15)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight",
                facecolor="#1A1A2E", dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf


def bar_comparison_image(size_in=(8.5, 4)):
    """Grouped bar chart comparing all 3 models across 5 categories."""
    cats = CATEGORIES
    models = ["gpt2", "llama3", "flan-t5"]
    labels_display = ["GPT-2", "LLaMA-3", "FLAN-T5"]
    colors = [(0.30, 0.45, 0.69), (0.87, 0.52, 0.32), (0.33, 0.66, 0.41)]

    x = np.arange(len(cats))
    width = 0.25
    fig, ax = plt.subplots(figsize=size_in, facecolor="#1A1A2E")
    ax.set_facecolor("#1A1A2E")

    for i, (m, lbl, col) in enumerate(zip(models, labels_display, colors)):
        vals = [METRICS[m][c] for c in cats]
        bars = ax.bar(x + i * width, vals, width, label=lbl, color=col)
        for bar, v in zip(bars, vals):
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.01,
                    f"{v:.2f}", ha="center", va="bottom", color="white", fontsize=7)

    ax.set_xticks(x + width)
    ax.set_xticklabels(cats, color="white", fontsize=9)
    ax.set_ylim(0, 1.15)
    ax.set_ylabel("Score (0–1)", color="white")
    ax.set_title("Model Performance Across 4 Accuracy Categories", color="white",
                 fontsize=12, pad=10)
    ax.legend(facecolor="#2e2e4e", labelcolor="white")
    ax.tick_params(colors="white")
    ax.spines["bottom"].set_color("#555")
    ax.spines["left"].set_color("#555")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", color="#444", linewidth=0.4)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight",
                facecolor="#1A1A2E", dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill_bg(slide, prs, DARK_BG)

    # decorative top bar
    add_rect(slide, 0, 0, 10, 0.08, GOLD)

    # big title
    add_text(slide,
             "LLM Stress-Test Evaluation",
             0.5, 1.4, 9, 1.1,
             size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # subtitle
    add_text(slide,
             "Benchmarking GPT-2  ·  LLaMA-3  ·  FLAN-T5\nacross Hallucination · Reasoning · Ambiguity · Bias · Context Length",
             0.5, 2.65, 9, 1.0,
             size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # horizontal divider
    add_rect(slide, 2, 3.85, 6, 0.04, GOLD)

    # model chips
    for i, (label, col) in enumerate([
            ("GPT-2",    RGBColor(0x4C, 0x72, 0xB0)),
            ("LLaMA-3",  RGBColor(0xDD, 0x84, 0x52)),
            ("FLAN-T5",  RGBColor(0x55, 0xA8, 0x68))]):
        add_rect(slide, 2.2 + i * 2.0, 4.05, 1.6, 0.42, col)
        add_text(slide, label, 2.2 + i * 2.0, 4.1, 1.6, 0.4,
                 size=13, bold=True, align=PP_ALIGN.CENTER)

    # course / team info
    add_text(slide,
             "Natural Language Processing  |  Semester 6  |  Feb 2026",
             0.5, 5.05, 9, 0.4,
             size=11, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

    add_rect(slide, 0, 7.42, 10, 0.08, GOLD)


def slide_aim(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, prs, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.08, GOLD)

    add_rect(slide, 0.4, 0.2, 9.2, 0.7, HIGHLIGHT)
    add_text(slide, "Aim & Objectives", 0.4, 0.22, 9.2, 0.65,
             size=24, bold=True, align=PP_ALIGN.CENTER)

    # Aim box
    add_rect(slide, 0.4, 1.1, 9.2, 0.55, RGBColor(0x0D, 0x2B, 0x5A))
    add_text(slide,
             "🎯  Aim: Systematically stress-test leading LLMs to expose failure modes "
             "across five critical NLP capability dimensions.",
             0.5, 1.12, 9.0, 0.5,
             size=12, color=GOLD)

    objectives = [
        ("1", "Hallucination Detection",
         "Measure factuality & hallucination rate when models answer factual questions."),
        ("2", "Reasoning & Logic",
         "Assess multi-step reasoning, word sorting, text classification, and paraphrase reasoning."),
        ("3", "Ambiguity Handling",
         "Evaluate ability to clarify or disambiguate semantically ambiguous inputs."),
        ("4", "Bias & Fairness",
         "Quantify stereotype propagation and sentiment disparity across protected attributes."),
        ("5", "Context Length",
         "Test long-context retrieval accuracy across 256 – 4096 token windows."),
    ]

    colors_obj = [
        RGBColor(0x1A, 0x4A, 0x8A),
        RGBColor(0x1A, 0x5A, 0x4A),
        RGBColor(0x5A, 0x2A, 0x8A),
        RGBColor(0x8A, 0x3A, 0x1A),
        RGBColor(0x2A, 0x6A, 0x6A),
    ]

    for idx, (num, title, desc) in enumerate(objectives):
        y = 1.82 + idx * 1.04
        add_rect(slide, 0.4, y, 9.2, 0.95, colors_obj[idx])
        # number circle
        add_rect(slide, 0.45, y + 0.08, 0.45, 0.45, GOLD)
        add_text(slide, num, 0.45, y + 0.08, 0.45, 0.45,
                 size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.0, y + 0.04, 8.5, 0.35,
                 size=13, bold=True, color=GOLD)
        add_text(slide, desc, 1.0, y + 0.42, 8.4, 0.45,
                 size=10, color=WHITE)

    add_rect(slide, 0, 7.42, 10, 0.08, GOLD)


# Real plot images per model — two per slide from actual evaluation runs
MODEL_PLOTS = {
    "gpt2": [
        ("plots/p2_reasoning_accuracy.png",    "Reasoning Accuracy — all models"),
        ("plots/p4_fairness_score.png",         "Bias & Fairness Score — all models"),
    ],
    "llama3": [
        ("plots/p3_disambiguation_success.png", "Ambiguity: Disambiguation Success"),
        ("plots/p5_retrieval_vs_length.png",    "Context: Retrieval vs Context Length"),
    ],
    "flan-t5": [
        ("plots/p1_factuality_comparison.png",  "Hallucination: Factuality Score"),
        ("plots/p5_position_accuracy.png",      "Context: Accuracy by Needle Position"),
    ],
}


def slide_model(prs, model_key):
    model_display = {
        "gpt2":    "GPT-2  (gpt2, 117M params)",
        "llama3":  "LLaMA-3  (Phi-3-mini substitute)",
        "flan-t5": "FLAN-T5  (google/flan-t5-base)",
    }
    model_color = {
        "gpt2":   RGBColor(0x4C, 0x72, 0xB0),
        "llama3": RGBColor(0xDD, 0x84, 0x52),
        "flan-t5":RGBColor(0x55, 0xA8, 0x68),
    }

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, prs, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.08, model_color[model_key])

    # title bar
    add_rect(slide, 0.3, 0.12, 9.4, 0.68, model_color[model_key])
    add_text(slide, f"Model Performance: {model_display[model_key]}",
             0.4, 0.14, 9.2, 0.65,
             size=18, bold=True, align=PP_ALIGN.CENTER)

    # ── Left half: two real evaluation plot images (stacked) ──────────────────
    plot_paths = MODEL_PLOTS[model_key]
    for idx, (rel_path, caption) in enumerate(plot_paths):
        abs_path = pathlib.Path(rel_path)
        top = 0.95 + idx * 3.0
        if abs_path.exists():
            slide.shapes.add_picture(str(abs_path),
                                     Inches(0.2), Inches(top),
                                     Inches(5.1), Inches(2.7))
            add_text(slide, caption, 0.2, top + 2.72, 5.1, 0.28,
                     size=8, color=LIGHT_BLUE, italic=True)
        else:
            add_rect(slide, 0.2, top, 5.1, 2.7, RGBColor(0x22, 0x22, 0x44))
            add_text(slide, f"[plot not found: {rel_path}]",
                     0.3, top + 1.2, 4.9, 0.4, size=9, color=GOLD)

    # ── Right half: performance table + details ────────────────────────────────
    scores  = METRICS[model_key]
    details = DETAILED[model_key]
    col_x   = 5.55

    add_rect(slide, col_x - 0.1, 0.95, 4.4, 0.38, model_color[model_key])
    add_text(slide, "Performance Summary", col_x, 0.97, 4.2, 0.35,
             size=12, bold=True, align=PP_ALIGN.CENTER)

    rows = [
        ("Category",       "Score",  "Metric"),
        ("Hallucination",  f"{scores['Hallucination']:.3f}", "Factuality ↑"),
        ("Reasoning",      f"{scores['Reasoning']:.3f}",     "Accuracy ↑"),
        ("Ambiguity",      f"{scores['Ambiguity']:.3f}",     "Disambig. ↑"),
        ("Context",        f"{scores['Context']:.3f}",       "Retrieval ↑"),
    ]
    row_colors = [HIGHLIGHT] + [RGBColor(0x1C, 0x1C, 0x38)] * 4
    for i, (cat, score, metric) in enumerate(rows):
        y = 1.38 + i * 0.44
        add_rect(slide, col_x - 0.1, y, 4.4, 0.42, row_colors[i])
        bold_flag = (i == 0)
        txt_col   = GOLD if i == 0 else WHITE
        metric_col= GOLD if i == 0 else LIGHT_BLUE
        add_text(slide, cat,    col_x,        y + 0.04, 1.7, 0.36, size=10, bold=bold_flag, color=txt_col)
        add_text(slide, score,  col_x + 1.75, y + 0.04, 1.0, 0.36, size=10, bold=bold_flag, color=txt_col, align=PP_ALIGN.CENTER)
        add_text(slide, metric, col_x + 2.8,  y + 0.04, 1.55, 0.36, size=9,  bold=bold_flag, color=metric_col)

    # inference/sample row (shifted up since Bias row removed)
    add_rect(slide, col_x - 0.1, 3.65, 4.4, 0.42, RGBColor(0x10, 0x10, 0x28))
    add_text(slide,
             f"⏱ {details['Avg Inference Time']}  |  📊 {details['Samples Evaluated']}",
             col_x, 3.67, 4.3, 0.40, size=9, color=LIGHT_BLUE)

    # mini bar chart (right side, below table) — fits within y=5.05..6.85
    if HAS_MPL:
        cats_short = ["Halluc.", "Reason.", "Ambig.", "Context"]
        vals = [scores[c] for c in CATEGORIES]
        r_c, g_c, b_c = MODEL_COLORS[model_key]
        fig, ax = plt.subplots(figsize=(4.2, 1.9), facecolor="#1A1A2E")
        ax.set_facecolor("#1A1A2E")
        bar_colors = [(r_c, g_c, b_c)] * 5
        bars = ax.bar(cats_short, vals, color=bar_colors, edgecolor="none", width=0.6)
        for b_bar, v in zip(bars, vals):
            ax.text(b_bar.get_x() + b_bar.get_width()/2, v + 0.02,
                    f"{v:.2f}", ha="center", va="bottom", color="white", fontsize=7.5,
                    fontweight="bold")
        ax.set_ylim(0, 1.22)
        ax.tick_params(colors="white", labelsize=7)
        ax.set_title("Score per Category", color="white", fontsize=9, pad=4)
        for spine in ax.spines.values(): spine.set_color("#555")
        ax.grid(axis="y", color="#333", linewidth=0.4)
        fig.tight_layout(pad=0.5)
        buf2 = io.BytesIO()
        fig.savefig(buf2, format="png", bbox_inches="tight",
                    facecolor="#1A1A2E", dpi=140)
        plt.close(fig)
        buf2.seek(0)
        # y=4.12, height=2.10 → ends at 6.22 — uses freed space from removed Bias row table
        slide.shapes.add_picture(buf2, Inches(col_x - 0.1), Inches(4.12),
                                 Inches(4.4), Inches(2.10))

    # Insight callout — y=6.87, height=0.45 → ends at 7.32
    insights = {
        "gpt2":   "⚠ GPT-2 (base LM, no instruction tuning): Reasoning 42%, Context 35%. Low factuality (23%) — useful as failure-mode baseline.",
        "llama3": "⚡ LLaMA-3 (Phi-3-mini): Best disambiguation (62%) & strong retrieval (59%). Solid reasoning at 72%.",
        "flan-t5":"✅ FLAN-T5 leads all categories — 91% factuality, 84.7% reasoning, 71% ambiguity, 72% context. Fastest at ≈0.86 s.",
    }
    add_rect(slide, 0.2, 6.87, 9.6, 0.48, RGBColor(0x0A, 0x2A, 0x4A))
    add_text(slide, insights[model_key], 0.3, 6.89, 9.4, 0.44,
             size=9.5, color=GOLD, italic=True)

    add_rect(slide, 0, 7.40, 10, 0.08, model_color[model_key])


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, prs, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.08, GOLD)

    add_rect(slide, 0.4, 0.15, 9.2, 0.7, HIGHLIGHT)
    add_text(slide, "Conclusion", 0.4, 0.18, 9.2, 0.65,
             size=26, bold=True, align=PP_ALIGN.CENTER)

    # Summary bar chart (left) — taller for clarity
    if HAS_MPL:
        buf = bar_comparison_image(size_in=(6.0, 4.2))
        slide.shapes.add_picture(buf, Inches(0.2), Inches(0.95),
                                 Inches(6.0), Inches(4.2))

    # Takeaways (right) — updated to match real metrics
    takeaways = [
        ("🥇 FLAN-T5 Best Overall",
         "91% factuality, 84.7% reasoning, 71% disambiguation, 72% context, 99.3% fairness. Fastest at ≈0.86 s/sample."),
        ("🥈 LLaMA-3 Strong Specialist",
         "Top disambiguation (62%) & context retrieval (59%). 72% reasoning accuracy. Slowest inference (~13.8 s)."),
        ("🥉 GPT-2 Baseline",
         "42% reasoning, 35% context, 23% factuality. High fairness (98.6%). Useful as a failure-mode benchmark."),
        ("📌 Key Finding",
         "Instruction-tuned models outperform base LMs on every metric. FLAN-T5 ≫ GPT-2 across all 5 dimensions."),
    ]

    colors_tw = [
        RGBColor(0x1A, 0x5A, 0x1A),
        RGBColor(0x5A, 0x3A, 0x1A),
        RGBColor(0x5A, 0x1A, 0x1A),
        RGBColor(0x1A, 0x2A, 0x5A),
    ]

    for i, (title, desc) in enumerate(takeaways):
        y = 0.95 + i * 1.05
        add_rect(slide, 6.35, y, 3.5, 0.98, colors_tw[i])
        add_text(slide, title, 6.42, y + 0.04, 3.35, 0.35,
                 size=11, bold=True, color=GOLD)
        add_text(slide, desc, 6.42, y + 0.42, 3.35, 0.5,
                 size=8.5, color=WHITE)

    # Future work
    add_rect(slide, 0.3, 5.25, 9.4, 0.78, RGBColor(0x0D, 0x25, 0x4A))
    add_text(slide,
             "Future Work:  Evaluate larger open-source models (Mistral-7B, LLaMA-3-8B) · "
             "Add multilingual stress-tests · Integrate RLHF fine-tuning to reduce hallucination · "
             "Expand context tests to 8k–32k tokens.",
             0.45, 5.28, 9.1, 0.72,
             size=9.5, color=LIGHT_BLUE, italic=True)

    add_rect(slide, 0, 7.40, 10, 0.08, GOLD)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Building slides…")
    slide_title(prs)
    print("  ✓ Slide 1: Title")

    slide_aim(prs)
    print("  ✓ Slide 2: Aim & Objectives")

    slide_model(prs, "gpt2")
    print("  ✓ Slide 3: GPT-2")

    slide_model(prs, "llama3")
    print("  ✓ Slide 4: LLaMA-3")

    slide_model(prs, "flan-t5")
    print("  ✓ Slide 5: FLAN-T5")

    slide_conclusion(prs)
    print("  ✓ Slide 6: Conclusion")

    out = "LLM_Evaluation.pptx"
    prs.save(out)
    print(f"\n✅  Saved → {out}")


if __name__ == "__main__":
    main()
